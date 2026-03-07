"""
Generate PBL PPT for Rapid Care Blood Distribution Network Platform.
Run: python create_pptx.py   (or: .venv/bin/python create_pptx.py after pip install python-pptx)
Output: PBL-Rapid-Care-Blood-Distribution.pptx
"""
from pptx import Presentation
from pptx.util import Pt

STUDENT_NAME = "Pudota Ashritha"
ROLL_NUMBER = "2520030592"
PROJECT_TITLE = "Rapid Care Blood Distribution Network Platform"

def add_title_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)

def add_content_slide(prs, title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    if points and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, pt in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = pt
            p.font.size = Pt(13)

def main():
    prs = Presentation()
    add_title_slide(prs, PROJECT_TITLE, [
        f"Name: {STUDENT_NAME}", f"Roll Number: {ROLL_NUMBER}",
        "PBL – FWD & DSA", "College: KLH", "[Add KLH Logo – Insert → Pictures]"
    ])
    add_content_slide(prs, "Course & Project Details", [
        f"Student: {STUDENT_NAME}", f"Roll No: {ROLL_NUMBER}",
        "Project: Rapid Care Blood Distribution Network Platform",
        "PBL – FWD (Web) & DSA", "College: KLH"
    ])
    add_content_slide(prs, "Abstract", [
        "Web-based application to manage blood distribution requests.",
        "Add requests (blood group, units, hospital); view list; search; sort.",
        "Tabs: Home, Login, Register, Donor, Stock. Built with HTML, CSS, JavaScript.",
        "DSA version: Stack (Undo), Queue (Process Next), Linear Search, Bubble Sort. Java module: LinkedList, Stack, Queue, HashMap."
    ])
    add_content_slide(prs, "Existing System", [
        "Blood distribution tracked on paper or spreadsheets.",
        "Hard to search, sort, process in order; no single view; no undo."
    ])
    add_content_slide(prs, "Proposed System", [
        "Web platform: add requests, list, search, sort.",
        "DSA: Undo (Stack), Process Next (Queue). Java: LinkedList, Stack, Queue, HashMap."
    ])
    add_content_slide(prs, "Key Features", [
        "Five tabs: Home, Login, Register, Donor, Stock.",
        "Add, search, sort, manage requests; Login/Register forms; Donor & Stock tables.",
        "DSA (JS): Undo, Process Next, Search, Sort. DSA (Java): LinkedList, Stack, Queue, Hashing."
    ])
    add_content_slide(prs, "Technologies Used", [
        "HTML – structure; CSS – styling; JavaScript – tabs, forms, list, DSA in browser.",
        "Java – LinkedList, Stack, Queue, HashMap (console program)."
    ])
    add_content_slide(prs, "Module 1: UI (HTML & CSS)", [
        "Header, nav (Home, Login, Register, Donor, Stock), forms, tables.",
        "Layout and styling with CSS; blood theme (red/gray)."
    ])
    add_content_slide(prs, "Module 2: Request Management (JavaScript)", [
        "Add, remove, display requests; array + DOM updates."
    ])
    add_content_slide(prs, "Module 3: Search & Sort", [
        "Linear search; bubble sort (DSA version)."
    ])
    add_content_slide(prs, "Module 4: DSA (JS + Java)", [
        "JS: Stack (Undo), Queue (Process Next), Search, Sort – console output.",
        "Java: LinkedList, Stack, Queue, HashMap – menu program in java-dsa/."
    ])
    add_content_slide(prs, "Screenshot: Home", ["[Insert screenshot: Home tab with requests]"])
    add_content_slide(prs, "Screenshot: Login / Register", ["[Insert screenshot: Login or Register form]"])
    add_content_slide(prs, "Screenshot: Donor / Stock", ["[Insert screenshot: Donor or Stock table]"])
    add_content_slide(prs, "Screenshot: DSA Console", ["[Insert screenshot: DSA version with F12 console]"])
    add_content_slide(prs, "Conclusion", [
        "Developed Rapid Care using HTML, CSS, JavaScript; tabs and request management.",
        "DSA in browser (Stack, Queue, Search, Sort) and in Java (LinkedList, Stack, Queue, HashMap)."
    ])
    add_content_slide(prs, "Future Scope", [
        "localStorage; filters; admin dashboard; backend; auth."
    ])
    add_content_slide(prs, "Geo-tagged: Demo 1", ["[Insert geo-tagged photo – Demo 1]"])
    add_content_slide(prs, "Geo-tagged: Demo 2", ["[Insert geo-tagged photo – Demo 2]"])
    add_content_slide(prs, "Geo-tagged: Demo 3", ["[Insert geo-tagged photo – Demo 3]"])
    add_title_slide(prs, "Thank You", ["Q&A", f"{STUDENT_NAME} | Roll No: {ROLL_NUMBER}"])
    add_content_slide(prs, "References", ["MDN, W3Schools, Course syllabus – FWD & DSA"])

    out = "PBL-Rapid-Care-Blood-Distribution.pptx"
    prs.save(out)
    print(f"Done! Saved: {out} | Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
